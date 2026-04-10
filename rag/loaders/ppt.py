import importlib
from tqdm import tqdm
from typing import Literal
from pptx import Presentation
from unstructured.partition.text import partition_text
from langchain_community.document_loaders import UnstructuredFileLoader

from rag.loaders.img import IMAGE_LOADER_METHODS


class PPTLoader(UnstructuredFileLoader):
    def __init__(self, file_path, method: Literal["rapidocr", "llmcaption", "blipcaption"]="rapidocr", **kwargs):
        super().__init__(file_path, **kwargs)
        module = importlib.import_module("rag.loaders.img")
        ImageLoader = getattr(module, IMAGE_LOADER_METHODS[method])
        self.image_loader = ImageLoader()
    def ppt2text(self, file_path):
        content = ""
        ppt = Presentation(file_path)
        with tqdm(total=len(ppt.slides), position=0, leave=True, desc="PPT Loading") as pbar:
            for slide in ppt.slides: # 循环每页幻灯片
                sorted_shapes = sorted( # 从上到下，从左到右排序
                    slide.shapes, key=lambda x: (x.top, x.left))
                for shape in sorted_shapes: # 循环每个形状块
                    if shape.has_text_frame:
                        content += shape.text + "\n\n"
                    elif shape.has_table:
                        table_content = ""
                        table = shape.table
                        for row in table.rows:
                            table_content += "| "
                            for cell in row.cells:
                                table_content += cell.text
                                table_content += " | "
                            table_content += "\n\n"
                        content += f"这里有一个表格内容为：\n\n{table_content}表格至此结束。\n\n"
                    elif shape.shape_type == 13:
                        image=shape.image
                        if image.size[0] > 10 and image.size[1] > 10:
                            if image_content := self.image_loader.img2text(image.blob):
                                image_content = image_content.replace("\n", " ")
                                content += f"这里有一张图片内容为：“{image_content}”。\n\n"
                pbar.update()
        return content
    def _get_elements(self):
        text = self.ppt2text(self.file_path)
        return partition_text(text=text, **self.unstructured_kwargs)
